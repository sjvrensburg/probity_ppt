#!/usr/bin/env python3
"""Post-render step: turn marker slides into branded card / stat-card rows.

Pandoc's pptx writer cannot draw boxes or lay out three columns (it drops the
third). This script runs after `quarto render`, scans each output .pptx, and
replaces the body of any slide whose content begins with a marker token:

    [[cards]]      three-card row    (label :: body)
    [[statcards]]  stat callout row  (label :: number :: description)

Data lines use `::` as the field separator. Prefix a label with `*` to make
that card the gold emphasis variant (one per row, the most important figure).

Authoring example (.qmd):

    ## Headline results

    [[statcards]]

    - Sample :: 6 :: Fiscal-year observations
    - Fit :: 0.86 :: R-squared, lag-1 GDP
    - *Multiplier :: R 14.9M :: Applied provision uplift

Quarto invokes this through the `post-render` hook in `_quarto.yml`, so a plain
`quarto render` applies it automatically. It can also be run by hand:

    python3 build/probity_cards.py deck.pptx
"""
import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Palette
NAVY = "0A325A"; BG_OFF = "F7F9FC"; ACCENT = "C8881F"
TEXT = "1F2937"; MUTE = "6B7280"; RULE = "D5DEE9"; WHITE = "FFFFFF"
FONT = "Calibri"
W = 13.333

def rgb(h): return RGBColor.from_string(h)

def _fill(sp, color):
    sp.fill.solid(); sp.fill.fore_color.rgb = rgb(color)

def _noline(sp): sp.line.fill.background()

def rect(slide, x, y, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    _fill(sp, color)
    if line:
        sp.line.color.rgb = rgb(line); sp.line.width = Pt(0.75)
    else:
        _noline(sp)
    return sp

def text(slide, s, x, y, w, h, size, color, *, bold=False, spacing=None,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = s
    f = r.font; f.size = Pt(size); f.bold = bold; f.name = FONT; f.color.rgb = rgb(color)
    if spacing is not None:
        r._r.get_or_add_rPr().set("spc", str(int(spacing * 100)))
    return tb

# ---- parsing -------------------------------------------------------------
def body_placeholder(slide):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1 and ph.has_text_frame:
            return ph
    # fall back: any content placeholder with text
    for ph in slide.placeholders:
        if ph.placeholder_format.idx not in (0, 10, 11, 12) and ph.has_text_frame:
            return ph
    return None

def parse(ph):
    paras = [p.text.strip() for p in ph.text_frame.paragraphs]
    paras = [t for t in paras if t]
    if not paras:
        return None, []
    marker = paras[0].lower()
    if marker not in ("[[cards]]", "[[statcards]]"):
        return None, []
    items = []
    for line in paras[1:]:
        fields = [f.strip() for f in line.split("::")]
        gold = False
        if fields and fields[0].startswith("*"):
            gold = True; fields[0] = fields[0][1:].strip()
        items.append((fields, gold))
    return marker.strip("[]"), items

def remove_shape(shape):
    el = shape._element; el.getparent().remove(el)

# ---- drawing -------------------------------------------------------------
def layout_x(n, gap=0.3, margin=0.5):
    total = W - 2 * margin
    cardW = (total - (n - 1) * gap) / n
    return margin, cardW, gap

def draw_cards(slide, items):
    n = len(items)
    if n == 0: return
    startX, cardW, gap = layout_x(n)
    y, cardH = 2.2, 3.6
    for i, (fields, gold) in enumerate(items):
        x = startX + i * (cardW + gap)
        label = fields[0] if fields else ""
        body = fields[1] if len(fields) > 1 else ""
        rect(slide, x, y, cardW, cardH, NAVY if gold else BG_OFF, line=RULE if not gold else None)
        rect(slide, x, y, 0.08, cardH, ACCENT if gold else NAVY)  # left accent edge
        text(slide, label, x + 0.35, y + 0.4, cardW - 0.6, 0.9, 20, ACCENT if gold else NAVY,
             bold=True)
        text(slide, body, x + 0.35, y + 1.35, cardW - 0.6, cardH - 1.6, 14,
             WHITE if gold else TEXT)

def draw_statcards(slide, items):
    n = len(items)
    if n == 0: return
    startX, cardW, gap = layout_x(n)
    y, cardH = 2.3, 3.0
    for i, (fields, gold) in enumerate(items):
        x = startX + i * (cardW + gap)
        label = fields[0] if fields else ""
        number = fields[1] if len(fields) > 1 else ""
        desc = fields[2] if len(fields) > 2 else ""
        rect(slide, x, y, cardW, cardH, NAVY if gold else WHITE, line=None if gold else RULE)
        if not gold:
            rect(slide, x, y, cardW, 0.12, NAVY)  # navy top edge
        text(slide, label.upper(), x + 0.3, y + 0.35, cardW - 0.6, 0.5, 12,
             ACCENT if gold else MUTE, bold=True, spacing=3)
        text(slide, number, x + 0.3, y + 0.85, cardW - 0.6, 1.1, 48,
             WHITE if gold else NAVY, bold=True)
        text(slide, desc, x + 0.3, y + 2.05, cardW - 0.6, 0.8, 12,
             WHITE if gold else TEXT, anchor=MSO_ANCHOR.TOP)

def process_file(path):
    prs = Presentation(path)
    touched = 0
    for slide in prs.slides:
        ph = body_placeholder(slide)
        if ph is None:
            continue
        kind, items = parse(ph)
        if not kind:
            continue
        remove_shape(ph)
        if kind == "cards":
            draw_cards(slide, items)
        else:
            draw_statcards(slide, items)
        touched += 1
    if touched:
        prs.save(path)
    print(f"probity_cards: {touched} card slide(s) in {os.path.basename(path)}")

def main():
    files = sys.argv[1:]
    if not files:
        env = os.environ.get("QUARTO_PROJECT_OUTPUT_FILES", "")
        files = [f for f in env.splitlines() if f.strip()]
    for f in files:
        if f.lower().endswith(".pptx") and os.path.exists(f):
            process_file(f)

if __name__ == "__main__":
    main()
