#!/usr/bin/env python3
"""Build the Probity-branded pandoc reference.pptx from the default reference.

Steps:
1. Rewrite the theme colour scheme (navy/gold palette) inside the pptx zip.
2. Open with python-pptx, set widescreen size, white master background.
3. Decorate slide layouts:
   - Title Slide / Section Header: navy fill, gold left stripe, white logo,
     white title, gold subtitle.
   - Content layouts: small navy logo, hairline rule, navy footer wordmark,
     navy titles.
"""
import re, shutil, zipfile, os
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
ASSETS = os.path.join(ROOT, "assets")
SRC = os.path.join(HERE, "default-reference.pptx")
BASE = os.path.join(HERE, "reference_base.pptx")
OUT = os.path.join(ROOT, "_extensions", "probity-pptx", "reference.pptx")

# ---- Palette -------------------------------------------------------------
NAVY      = "0A325A"
NAVY_DEEP = "062340"
BLUE_MED  = "4A7BA8"
BLUE_LITE = "8BABCB"
BLUE_TINT = "E8EEF5"
BG_OFF    = "F7F9FC"
ACCENT    = "C8881F"   # gold
TEXT      = "1F2937"
MUTE      = "6B7280"
RULE      = "D5DEE9"
WHITE     = "FFFFFF"

# Theme colour map (Office roles -> Probity colours)
CLR_SCHEME = f"""<a:clrScheme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="Probity">\
<a:dk1><a:srgbClr val="{TEXT}"/></a:dk1>\
<a:lt1><a:srgbClr val="{WHITE}"/></a:lt1>\
<a:dk2><a:srgbClr val="{NAVY}"/></a:dk2>\
<a:lt2><a:srgbClr val="{BG_OFF}"/></a:lt2>\
<a:accent1><a:srgbClr val="{NAVY}"/></a:accent1>\
<a:accent2><a:srgbClr val="{ACCENT}"/></a:accent2>\
<a:accent3><a:srgbClr val="{BLUE_MED}"/></a:accent3>\
<a:accent4><a:srgbClr val="{BLUE_LITE}"/></a:accent4>\
<a:accent5><a:srgbClr val="{BLUE_TINT}"/></a:accent5>\
<a:accent6><a:srgbClr val="{MUTE}"/></a:accent6>\
<a:hlink><a:srgbClr val="{BLUE_MED}"/></a:hlink>\
<a:folHlink><a:srgbClr val="{MUTE}"/></a:folHlink>\
</a:clrScheme>"""

# ---- Step 1: rewrite theme inside the zip --------------------------------
def rewrite_theme():
    shutil.copy(SRC, BASE)
    tmp = BASE + ".tmp"
    with zipfile.ZipFile(BASE) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.namelist():
            data = zin.read(item)
            if item == "ppt/theme/theme1.xml":
                txt = data.decode("utf-8")
                txt = re.sub(r"<a:clrScheme.*?</a:clrScheme>", CLR_SCHEME, txt, flags=re.S)
                # major/minor fonts -> Calibri (already Calibri latin, ensure)
                txt = txt.replace('name="Office"><a:majorFont><a:latin typeface="Calibri"',
                                  'name="Probity"><a:majorFont><a:latin typeface="Calibri"')
                data = txt.encode("utf-8")
            zout.writestr(item, data)
    os.replace(tmp, BASE)

# ---- low-level XML shape helpers (work on master & layout spTree) ---------
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
P = "http://schemas.openxmlformats.org/presentationml/2006/main"
R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
EMU = 914400

def _emu(v): return str(int(round(v * EMU)))

def rgb(h): return RGBColor.from_string(h)

_uid = [100]
def _nextid():
    _uid[0] += 1
    return _uid[0]

def _sptree(obj):
    return obj._element.find(qn("p:cSld")).find(qn("p:spTree"))

def add_rect(obj, x, y, w, h, color):
    nid = _nextid()
    xml = f"""<p:sp xmlns:p="{P}" xmlns:a="{A}">
  <p:nvSpPr><p:cNvPr id="{nid}" name="Rect {nid}"/><p:cNvSpPr/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{_emu(x)}" y="{_emu(y)}"/><a:ext cx="{_emu(w)}" cy="{_emu(h)}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
    <a:ln><a:noFill/></a:ln>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p/></p:txBody>
</p:sp>"""
    _sptree(obj).append(etree.fromstring(xml))

def add_text(obj, s, x, y, w, h, size, color, *, bold=False, italic=False,
             align="l", font="Calibri", spacing=None, anchor="t"):
    nid = _nextid()
    spc = f' spc="{int(spacing*100)}"' if spacing is not None else ""
    xml = f"""<p:sp xmlns:p="{P}" xmlns:a="{A}">
  <p:nvSpPr><p:cNvPr id="{nid}" name="Txt {nid}"/><p:cNvSpPr txBox="1"/><p:nvPr/></p:nvSpPr>
  <p:spPr>
    <a:xfrm><a:off x="{_emu(x)}" y="{_emu(y)}"/><a:ext cx="{_emu(w)}" cy="{_emu(h)}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="square" lIns="0" tIns="0" rIns="0" bIns="0" anchor="{anchor}"><a:spAutoFit/></a:bodyPr>
    <a:lstStyle/>
    <a:p><a:pPr algn="{align}"/>
      <a:r><a:rPr lang="en-GB" sz="{int(size*100)}" b="{1 if bold else 0}" i="{1 if italic else 0}"{spc}>
        <a:solidFill><a:srgbClr val="{color}"/></a:solidFill>
        <a:latin typeface="{font}"/></a:rPr><a:t>{s}</a:t></a:r></a:p>
  </p:txBody>
</p:sp>"""
    _sptree(obj).append(etree.fromstring(xml))

def add_logo(obj, path, x, y, w, h):
    image_part, rId = obj.part.get_or_add_image_part(path)
    nid = _nextid()
    xml = f"""<p:pic xmlns:p="{P}" xmlns:a="{A}" xmlns:r="{R}">
  <p:nvPicPr><p:cNvPr id="{nid}" name="Logo {nid}"/>
    <p:cNvPicPr><a:picLocks noChangeAspect="1"/></p:cNvPicPr><p:nvPr/></p:nvPicPr>
  <p:blipFill><a:blip r:embed="{rId}"/><a:stretch><a:fillRect/></a:stretch></p:blipFill>
  <p:spPr>
    <a:xfrm><a:off x="{_emu(x)}" y="{_emu(y)}"/><a:ext cx="{_emu(w)}" cy="{_emu(h)}"/></a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
</p:pic>"""
    _sptree(obj).append(etree.fromstring(xml))

def style_placeholder_text(ph, color, size=None, bold=None, font="Calibri", align="l"):
    """Set default run props on a placeholder's first level so inherited text
    picks up the colour/size."""
    txBody = ph._element.find(qn("p:txBody"))
    if txBody is None:
        return
    # ensure a lstStyle with lvl1pPr/defRPr
    lstStyle = txBody.find(qn("a:lstStyle"))
    if lstStyle is None:
        lstStyle = etree.SubElement(txBody, qn("a:lstStyle"))
        txBody.remove(lstStyle); txBody.insert(1, lstStyle)
    lvl1 = lstStyle.find(qn("a:lvl1pPr"))
    if lvl1 is None:
        lvl1 = etree.SubElement(lstStyle, qn("a:lvl1pPr"))
    if align is not None:
        lvl1.set("algn", align)
    defRPr = lvl1.find(qn("a:defRPr"))
    if defRPr is None:
        defRPr = etree.SubElement(lvl1, qn("a:defRPr"))
        lvl1.remove(defRPr); lvl1.insert(0, defRPr)
    if size is not None:
        defRPr.set("sz", str(int(size * 100)))
    if bold is not None:
        defRPr.set("b", "1" if bold else "0")
    # font
    for tag in ("a:latin",):
        el = defRPr.find(qn(tag))
        if el is None:
            el = etree.SubElement(defRPr, qn(tag))
        el.set("typeface", font)
    # colour
    for old in defRPr.findall(qn("a:solidFill")):
        defRPr.remove(old)
    sf = etree.Element(qn("a:solidFill"))
    sc = etree.SubElement(sf, qn("a:srgbClr")); sc.set("val", color)
    defRPr.insert(0, sf)

def set_bg(obj, color):
    """Set solid background fill on a master/layout."""
    cSld = obj._element.find(qn("p:cSld"))
    old = cSld.find(qn("p:bg"))
    if old is not None:
        cSld.remove(old)
    bg = etree.Element(qn("p:bg"))
    bgPr = etree.SubElement(bg, qn("p:bgPr"))
    fill = etree.SubElement(bgPr, qn("a:solidFill"))
    sc = etree.SubElement(fill, qn("a:srgbClr")); sc.set("val", color)
    etree.SubElement(bgPr, qn("a:effectLst"))
    cSld.insert(0, bg)

# ---- Step 2/3 ------------------------------------------------------------
W, H = 13.333, 7.5
LOGO_WHITE = os.path.join(ASSETS, "logo_white.png")
LOGO_NAVY  = os.path.join(ASSETS, "logo_navy_small.png")

DARK_LAYOUTS = {"Title Slide", "Section Header"}

def decorate():
    prs = Presentation(BASE)
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    master = prs.slide_masters[0]
    set_bg(master, WHITE)
    # pandoc derives content geometry from the MASTER placeholders: align them
    # to the branded content chrome (title below the hairline rule, body lower).
    for ph in master.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:      # title
            ph.left, ph.top = Inches(0.5), Inches(1.12)
            ph.width, ph.height = Inches(W - 1.0), Inches(0.7)
            style_placeholder_text(ph, NAVY, size=30, bold=True)
        elif idx == 1:    # body / content
            ph.left, ph.top = Inches(0.5), Inches(1.98)
            ph.width, ph.height = Inches(W - 1.0), Inches(4.7)
        elif idx in (2, 3, 4):  # date / footer / slide number
            ph.top = Inches(7.05); ph.height = Inches(0.3)

    for layout in master.slide_layouts:
        name = layout.name
        if name in DARK_LAYOUTS:
            set_bg(layout, NAVY)
            # gold left stripe
            add_rect(layout, 0, 0, 0.18, H, ACCENT)
            # white logo top-left
            add_logo(layout, LOGO_WHITE, 0.7, 0.6, 2.4, 0.744)
            # title + subtitle styling
            for ph in layout.placeholders:
                t = ph.placeholder_format.type
                if ph.placeholder_format.idx == 0:  # title
                    ph.left, ph.top = Inches(0.7), Inches(2.4)
                    ph.width, ph.height = Inches(W - 1.4), Inches(1.6)
                    style_placeholder_text(ph, WHITE, size=54, bold=True)
                elif ph.placeholder_format.idx == 1:  # subtitle / body
                    ph.left, ph.top = Inches(0.7), Inches(4.1)
                    ph.width, ph.height = Inches(W - 1.4), Inches(1.0)
                    style_placeholder_text(ph, ACCENT, size=22, bold=False)
                elif ph.placeholder_format.idx in (10, 11, 12):  # date/footer/page
                    ph.left, ph.top = Inches(0.7), Inches(6.95)
                    ph.height = Inches(0.3)
                    style_placeholder_text(ph, BLUE_LITE, size=10, bold=False)
            # gold anchor rule below subtitle
            add_rect(layout, 0.7, 5.5, 2.2, 0.05, ACCENT)
        else:
            # content layouts: white bg inherited from master
            # small navy logo + hairline + footer wordmark
            add_logo(layout, LOGO_NAVY, 0.5, 0.32, 1.2, 0.372)
            add_rect(layout, 0.5, 0.92, W - 1.0, 0.02, RULE)
            add_text(layout, "Probity Data Analytics", 0.5, H - 0.42, 5, 0.3,
                     9, NAVY, bold=True)
            # move + style title placeholder (navy)
            for ph in layout.placeholders:
                if ph.placeholder_format.idx == 0:
                    ph.left, ph.top = Inches(0.5), Inches(1.12)
                    ph.width, ph.height = Inches(W - 1.0), Inches(0.7)
                    style_placeholder_text(ph, NAVY, size=30, bold=True)
                # nudge body/content placeholders down below the rule
                elif t_is_content(ph):
                    if ph.top < Inches(1.95):
                        ph.top = Inches(1.95)
    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print("wrote", OUT)

def t_is_content(ph):
    return ph.placeholder_format.idx not in (0, 10, 11, 12)

if __name__ == "__main__":
    rewrite_theme()
    decorate()
